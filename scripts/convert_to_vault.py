import os
import shutil
import pymupdf4llm
import mammoth
from pptx import Presentation
import uuid

INPUT_DIR = "input_docs"
OUTPUT_DIR = "site/content"
ASSETS_DIR = os.path.join(OUTPUT_DIR, "assets")

def process_files():
    # Cria as pastas do site e de imagens se não existirem
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(ASSETS_DIR, exist_ok=True)

    for filename in os.listdir(INPUT_DIR):
        input_path = os.path.join(INPUT_DIR, filename)
        
        if not os.path.isfile(input_path):
            continue
            
        name, ext = os.path.splitext(filename)
        ext = ext.lower()
        output_path = os.path.join(OUTPUT_DIR, f"{name}.md")

        try:
            if ext == '.md':
                shutil.copy2(input_path, output_path)
                print(f"Copiado direto: {filename}")
                
            elif ext == '.pdf':
                # PDF: Extrai o texto e salva as imagens na pasta assets
                md_text = pymupdf4llm.to_markdown(
                    input_path,
                    write_images=True,
                    image_path=ASSETS_DIR,
                    image_format="png"
                )
                
                # Limpa o caminho das imagens no texto para o site conseguir encontrar
                md_text = md_text.replace("site/content/assets/", "assets/")
                md_text = md_text.replace("site\\content\\assets\\", "assets/")
                
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(md_text)
                print(f"Convertido de PDF com imagens: {filename}")
                
            elif ext == '.docx':
                # DOCX: Função para recortar imagens do Word e salvar na pasta assets
                def handle_image(image):
                    ext_img = image.content_type.split("/")[1]
                    # Cria um nome único para não misturar fotos de arquivos diferentes
                    img_name = f"{name.replace(' ', '_')}_{uuid.uuid4().hex[:6]}.{ext_img}"
                    img_path = os.path.join(ASSETS_DIR, img_name)
                    
                    with image.open() as img_bytes:
                        with open(img_path, "wb") as f:
                            f.write(img_bytes.read())
                            
                    # Retorna o link para o texto apontando para a pasta certa
                    return {"src": f"assets/{img_name}"}

                with open(input_path, "rb") as docx_file:
                    result = mammoth.convert_to_markdown(
                        docx_file,
                        convert_image=mammoth.images.inline(handle_image)
                    )
                    text = f"# {name}\n\n{result.value}"
                    
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                print(f"Convertido de DOCX com imagens: {filename}")
                
            elif ext == '.pptx':
                # O NOVO MOTOR PARA POWERPOINT: Lê slide por slide (Textos e Imagens)!
                prs = Presentation(input_path)
                md_text = f"# {name}\n\n"
                
                for i, slide in enumerate(prs.slides):
                    md_text += f"## Slide {i+1}\n\n" # Cria um título para cada slide
                    
                    for shape in slide.shapes:
                        # 1. Tenta extrair texto
                        if hasattr(shape, "text") and shape.text.strip():
                            md_text += shape.text + "\n\n"
                            
                        # 2. Tenta extrair imagem
                        if hasattr(shape, "image"):
                            image = shape.image
                            ext_img = image.ext # Pega a extensão da imagem (png, jpg, etc)
                            
                            # Cria um nome único para a imagem ir para a pasta assets
                            img_name = f"{name.replace(' ', '_')}_slide{i+1}_{uuid.uuid4().hex[:6]}.{ext_img}"
                            img_path = os.path.join(ASSETS_DIR, img_name)
                            
                            # Salva a imagem fisicamente
                            with open(img_path, "wb") as img_file:
                                img_file.write(image.blob)
                                
                            # Escreve o link da imagem no arquivo Markdown
                            md_text += f"![Imagem Slide {i+1}](assets/{img_name})\n\n"
                            
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(md_text)
                print(f"Convertido de PPTX com imagens: {filename}")
            else:
                print(f"Formato ignorado: {filename}")
                
        except Exception as e:
            print(f"Erro ao processar {filename}: {e}")

if __name__ == "__main__":
    process_files()